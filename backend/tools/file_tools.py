"""File reading tools for uploaded files and project files."""

from pathlib import Path

from langchain_core.tools import tool

MAX_TEXT_SIZE = 50 * 1024  # 50KB 文本截断限制
MAX_PDF_PAGES = 20  # PDF 最多读取 20 页
MAX_EXCEL_ROWS = 500  # Excel 最多读取 500 行


@tool
async def read_file(
    file_id: str,
    project_id: str | None = None,
    section: str | None = None,
) -> str:
    """读取文件内容 (统一工具，支持上传文件和项目文件)。

    这是读取文件的首选工具，可以处理两种类型的文件：
    - 用户上传文件：只需提供 file_id
    - 项目文件：需要同时提供 file_id 和 project_id

    支持的文件类型：
    - 文本文件：txt, md, json, csv, 代码文件
    - PDF 文件：提取文本内容
    - Word 文件：docx
    - Excel 文件：xlsx, xls
    - PowerPoint 文件：pptx

    Args:
        file_id: 文件 ID
        project_id: 项目 ID (项目文件需要，上传文件可省略)
        section: 可选，指定读取的部分
            - Excel: sheet 名称（如 "Sheet1"）
            - PDF: 页码范围（如 "1-5" 或 "3"）
            - 其他格式忽略此参数

    Returns:
        文件内容的文本形式
    """
    if project_id:
        # 项目文件路径 - 异步读取
        return await _read_project_file_async(file_id, project_id, section)
    else:
        # 上传文件路径 - 本地文件读取保持同步
        return _read_uploaded_file_impl(file_id, section)


async def _read_project_file_async(
    file_id: str,
    project_id: str,
    section: str | None = None,
) -> str:
    """异步读取项目文件"""
    try:
        storage_path = await _get_project_file_path(file_id, project_id)
    except Exception as e:
        return f"错误：获取文件路径失败 - {e}"

    if not storage_path:
        return f"错误：找不到项目文件 ID {file_id}"

    file_path = Path(storage_path)
    if not file_path.exists():
        return f"错误：文件不存在 {storage_path}"

    # 文件内容读取仍然是同步的（本地 I/O）
    return _read_file_content(file_path, section)


def _read_uploaded_file_impl(file_id: str, section: str | None = None) -> str:
    """读取上传文件的实现"""
    file_dir = Path(f"/tmp/sunnyagent_files/{file_id}")
    if not file_dir.exists():
        return f"错误：找不到文件 ID {file_id}"

    files = list(file_dir.iterdir())
    if not files:
        return f"错误：文件 ID {file_id} 目录为空"

    file_path = files[0]
    return _read_file_content(file_path, section)


def _read_file_content(file_path: Path, section: str | None = None) -> str:
    """读取文件内容的通用实现"""
    ext = file_path.suffix.lower()

    # PDF 文件 - 支持 section 参数指定页码
    if ext == ".pdf":
        return _read_pdf_with_section(file_path, section)

    # Excel 文件 - 支持 section 参数指定 sheet
    if ext in {".xlsx", ".xls"}:
        return _read_excel_with_section(file_path, section)

    # Word 文件 (docx)
    if ext == ".docx":
        try:
            from docx import Document

            doc = Document(str(file_path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            content = "\n\n".join(paragraphs)
            if len(content) > MAX_TEXT_SIZE:
                return (
                    content[:MAX_TEXT_SIZE]
                    + f"\n\n[... 内容已截断，原文件共 {len(content)} 字符 ...]"
                )
            return content
        except Exception as e:
            return f"读取 Word 文件失败：{e}"

    # 旧版 Word 文件 (doc) - 不支持直接读取
    if ext == ".doc":
        return (
            f"文件 '{file_path.name}' 是旧版 Word 格式 (.doc)，无法直接读取。\n"
            "建议：请将文件另存为 .docx 格式后重新上传。"
        )

    # PowerPoint 文件 (pptx)
    if ext == ".pptx":
        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))
            slides_text = []

            for i, slide in enumerate(prs.slides, 1):
                slide_content = [f"[Slide {i}]"]
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = getattr(shape, "text", "")
                        if text and text.strip():
                            slide_content.append(text)
                slides_text.append("\n".join(slide_content))

            return "\n\n".join(slides_text)
        except Exception as e:
            return f"读取 PowerPoint 文件失败：{e}"

    # 旧版 PowerPoint 文件 (ppt) - 不支持直接读取
    if ext == ".ppt":
        return (
            f"文件 '{file_path.name}' 是旧版 PowerPoint 格式 (.ppt)，无法直接读取。\n"
            "建议：请将文件另存为 .pptx 格式后重新上传。"
        )

    # 文本文件
    text_extensions = {".txt", ".md", ".json", ".csv"}
    if ext in text_extensions:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            if len(content) > MAX_TEXT_SIZE:
                return (
                    content[:MAX_TEXT_SIZE]
                    + f"\n\n[... 内容已截断，原文件共 {len(content)} 字符 ...]"
                )
            return content
        except Exception as e:
            return f"读取文件失败：{e}"

    # 尝试使用 FileExtractor 处理其他文件类型
    try:
        from backend.services.file_extractor import FileExtractor

        extractor = FileExtractor()
        result = extractor.extract(file_path)
        return result.raw_text
    except Exception as e:
        return f"读取文件失败：{e}（不支持的文件类型：{ext}）"


# Legacy function - kept for backwards compatibility
@tool
async def read_project_file(
    file_id: str,
    project_id: str,
    section: str | None = None,
) -> str:
    """读取项目文件的完整内容或指定部分。

    当项目文件内容被摘要时（因文件过大），使用此工具获取完整内容。

    支持的文件类型：
    - 文本文件：txt, md, json, csv, 代码文件
    - PDF 文件：提取文本内容
    - Word 文件：docx
    - Excel 文件：xlsx, xls
    - PowerPoint 文件：pptx

    Args:
        file_id: 项目文件 ID
        project_id: 项目 ID
        section: 可选，指定读取的部分
            - Excel: sheet 名称（如 "Sheet1"）
            - PDF: 页码范围（如 "1-5" 或 "3"）
            - 其他格式忽略此参数

    Returns:
        文件内容的文本形式
    """
    # 使用统一的异步实现
    return await _read_project_file_async(file_id, project_id, section)


async def _get_project_file_path(file_id: str, project_id: str) -> str | None:
    """从数据库获取项目文件的存储路径"""
    from backend.db import get_pool

    pool = await get_pool()
    async with pool.acquire() as conn:
        row = await conn.fetchrow(
            """
            SELECT pf.storage_path
            FROM project_files pf
            JOIN projects p ON pf.project_id = p.id
            WHERE pf.file_id = $1
              AND p.id = $2
              AND p.is_deleted = FALSE
            """,
            file_id,
            project_id,
        )
        return row["storage_path"] if row else None


def _read_pdf_with_section(file_path: Path, section: str | None) -> str:
    """读取 PDF 文件，支持指定页码范围"""
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(file_path))
        total_pages = len(reader.pages)

        # 解析页码范围
        start_page, end_page = 0, total_pages
        if section:
            try:
                if "-" in section:
                    parts = section.split("-")
                    start_page = max(0, int(parts[0]) - 1)
                    end_page = min(total_pages, int(parts[1]))
                else:
                    page_num = int(section) - 1
                    start_page = max(0, page_num)
                    end_page = min(total_pages, page_num + 1)
            except ValueError:
                pass  # 无效的页码格式，读取全部

        pages_text = []
        for i in range(start_page, end_page):
            text = reader.pages[i].extract_text() or ""
            pages_text.append(f"[Page {i + 1}]\n{text}")

        result = "\n\n".join(pages_text)
        if end_page < total_pages:
            result += f"\n\n[提示：共 {total_pages} 页，当前显示第 {start_page + 1}-{end_page} 页]"

        return result
    except Exception as e:
        return f"读取 PDF 失败：{e}"


def _read_excel_with_section(file_path: Path, section: str | None) -> str:
    """读取 Excel 文件，支持指定 sheet"""
    try:
        from openpyxl import load_workbook

        wb = load_workbook(str(file_path), read_only=True, data_only=True)
        result_parts = []

        # 确定要读取的 sheet
        if section and section in wb.sheetnames:
            sheets_to_read = [section]
        else:
            sheets_to_read = wb.sheetnames

        for sheet_name in sheets_to_read:
            sheet = wb[sheet_name]
            result_parts.append(f"=== Sheet: {sheet_name} ===")

            rows = []
            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                if row_count >= MAX_EXCEL_ROWS:
                    result_parts.append(f"\n[... 仅显示前 {MAX_EXCEL_ROWS} 行 ...]")
                    break
                row_str = "\t".join(
                    str(cell) if cell is not None else "" for cell in row
                )
                if row_str.strip():
                    rows.append(row_str)
                    row_count += 1

            result_parts.append("\n".join(rows))

        wb.close()
        return "\n\n".join(result_parts)
    except Exception as e:
        return f"读取 Excel 文件失败：{e}"


@tool
def read_uploaded_file(file_id: str) -> str:
    """读取用户上传的文件内容。

    支持的文件类型：
    - 文本文件：txt, md, json, csv（直接读取，超过50KB截断）
    - PDF 文件：提取文本内容（最多读取前20页）
    - Word 文件：docx（提取文本内容）
    - Excel 文件：xlsx, xls（提取表格数据，最多500行）
    - PowerPoint 文件：pptx（提取幻灯片文本）

    Args:
        file_id: 文件 ID（从用户消息的附件信息中获取）

    Returns:
        文件内容的文本形式
    """
    file_dir = Path(f"/tmp/sunnyagent_files/{file_id}")
    if not file_dir.exists():
        return f"错误：找不到文件 ID {file_id}"

    files = list(file_dir.iterdir())
    if not files:
        return f"错误：文件 ID {file_id} 目录为空"

    file_path = files[0]
    ext = file_path.suffix.lower()

    # 文本文件
    if ext in {".txt", ".md", ".json", ".csv"}:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            if len(content) > MAX_TEXT_SIZE:
                return (
                    content[:MAX_TEXT_SIZE]
                    + f"\n\n[... 内容已截断，原文件共 {len(content)} 字符 ...]"
                )
            return content
        except Exception as e:
            return f"读取文件失败：{e}"

    # PDF 文件
    if ext == ".pdf":
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(file_path))
            total_pages = len(reader.pages)
            pages_to_read = min(total_pages, MAX_PDF_PAGES)

            pages_text = []
            for i in range(pages_to_read):
                text = reader.pages[i].extract_text() or ""
                pages_text.append(f"[Page {i+1}]\n{text}")

            result = "\n\n".join(pages_text)
            if total_pages > MAX_PDF_PAGES:
                result += f"\n\n[... 仅显示前 {MAX_PDF_PAGES} 页，原文件共 {total_pages} 页 ...]"
            return result
        except Exception as e:
            return f"读取 PDF 失败：{e}"

    # Word 文件 (docx)
    if ext == ".docx":
        try:
            from docx import Document

            doc = Document(str(file_path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            content = "\n\n".join(paragraphs)
            if len(content) > MAX_TEXT_SIZE:
                return (
                    content[:MAX_TEXT_SIZE]
                    + f"\n\n[... 内容已截断，原文件共 {len(content)} 字符 ...]"
                )
            return content
        except Exception as e:
            return f"读取 Word 文件失败：{e}"

    # 旧版 Word 文件 (doc) - 不支持直接读取
    if ext == ".doc":
        return (
            f"文件 '{file_path.name}' 是旧版 Word 格式 (.doc)，无法直接读取。\n"
            "建议：请将文件另存为 .docx 格式后重新上传，或使用 activate_skill('docx') 获取处理指南。"
        )

    # Excel 文件 (xlsx, xls)
    if ext in {".xlsx", ".xls"}:
        try:
            from openpyxl import load_workbook

            wb = load_workbook(str(file_path), read_only=True, data_only=True)
            result_parts = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                result_parts.append(f"=== Sheet: {sheet_name} ===")

                rows = []
                row_count = 0
                for row in sheet.iter_rows(values_only=True):
                    if row_count >= MAX_EXCEL_ROWS:
                        result_parts.append(f"\n[... 仅显示前 {MAX_EXCEL_ROWS} 行 ...]")
                        break
                    # Convert row to string, handling None values
                    row_str = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    if row_str.strip():  # Skip empty rows
                        rows.append(row_str)
                        row_count += 1

                result_parts.append("\n".join(rows))

            wb.close()
            return "\n\n".join(result_parts)
        except Exception as e:
            return f"读取 Excel 文件失败：{e}"

    # PowerPoint 文件 (pptx)
    if ext == ".pptx":
        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))
            slides_text = []

            for i, slide in enumerate(prs.slides, 1):
                slide_content = [f"[Slide {i}]"]
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = getattr(shape, "text", "")
                        if text and text.strip():
                            slide_content.append(text)
                slides_text.append("\n".join(slide_content))

            return "\n\n".join(slides_text)
        except Exception as e:
            return f"读取 PowerPoint 文件失败：{e}"

    # 旧版 PowerPoint 文件 (ppt) - 不支持直接读取
    if ext == ".ppt":
        return (
            f"文件 '{file_path.name}' 是旧版 PowerPoint 格式 (.ppt)，无法直接读取。\n"
            "建议：请将文件另存为 .pptx 格式后重新上传。"
        )

    return f"不支持的文件类型：{ext}"
