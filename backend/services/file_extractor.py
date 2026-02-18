"""
File content extraction service.

Provides unified extraction for various file types with structured output.
"""

import json
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from pathlib import Path


@dataclass
class ExtractedContent:
    """提取的文件内容"""

    raw_text: str  # 原始文本
    structured_data: dict | None  # 结构化数据（Excel 表格等）
    metadata: dict = field(default_factory=dict)  # 文件元信息
    estimated_tokens: int = 0  # 估算 token 数

    def __post_init__(self):
        if self.estimated_tokens == 0:
            # 简单估算: 1 token ≈ 3 字符（混合中英文）
            self.estimated_tokens = len(self.raw_text) // 3


class BaseExtractor(ABC):
    """文件提取器基类"""

    @abstractmethod
    def extract(self, file_path: Path) -> ExtractedContent:
        """提取文件内容"""
        pass

    def _estimate_tokens(self, text: str) -> int:
        """估算 token 数"""
        return len(text) // 3


class TextExtractor(BaseExtractor):
    """文本文件提取器 (txt, md, json, csv, 代码文件)"""

    MAX_SIZE = 500 * 1024  # 500KB 限制

    def extract(self, file_path: Path) -> ExtractedContent:
        try:
            # 检查文件大小
            file_size = file_path.stat().st_size
            if file_size > self.MAX_SIZE:
                # 大文件只读取前 500KB
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read(self.MAX_SIZE)
                    content += "\n\n... [文件内容过长，已截断] ..."
            else:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()

            return ExtractedContent(
                raw_text=content,
                structured_data=None,
                metadata={
                    "format": "text",
                    "suffix": file_path.suffix.lower(),
                    "size_bytes": file_size,
                    "truncated": file_size > self.MAX_SIZE,
                },
            )
        except Exception as e:
            return ExtractedContent(
                raw_text=f"[文件读取失败: {e}]",
                structured_data=None,
                metadata={"error": str(e)},
            )


class CsvExtractor(BaseExtractor):
    """CSV 文件提取器"""

    MAX_ROWS = 500  # 最大行数

    def extract(self, file_path: Path) -> ExtractedContent:
        import csv

        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                reader = csv.reader(f)
                rows = list(reader)

            if not rows:
                return ExtractedContent(
                    raw_text="[空 CSV 文件]",
                    structured_data={"type": "csv", "rows": []},
                    metadata={"format": "csv", "row_count": 0},
                )

            headers = rows[0]
            data_rows = rows[1 : self.MAX_ROWS + 1]
            truncated = len(rows) > self.MAX_ROWS + 1

            structured = {
                "type": "csv",
                "headers": headers,
                "row_count": len(rows) - 1,
                "sample_rows": [
                    {headers[i]: cell for i, cell in enumerate(row) if i < len(headers)}
                    for row in data_rows
                ],
                "truncated": truncated,
            }

            # 生成可读文本
            text_lines = [",".join(headers)]
            for row in data_rows[:50]:  # 文本只显示前 50 行
                text_lines.append(",".join(row))
            if truncated:
                text_lines.append(f"... [共 {len(rows) - 1} 行，已截断显示] ...")
            text = "\n".join(text_lines)

            return ExtractedContent(
                raw_text=text,
                structured_data=structured,
                metadata={
                    "format": "csv",
                    "row_count": len(rows) - 1,
                    "column_count": len(headers),
                    "truncated": truncated,
                },
            )
        except Exception as e:
            return ExtractedContent(
                raw_text=f"[CSV 读取失败: {e}]",
                structured_data=None,
                metadata={"error": str(e)},
            )


class ExcelExtractor(BaseExtractor):
    """Excel 文件提取器 (.xlsx/.xls)"""

    MAX_SAMPLE_ROWS = 20  # 采样行数
    MAX_TOTAL_ROWS = 500  # 最大处理行数

    def extract(self, file_path: Path) -> ExtractedContent:
        try:
            from openpyxl import load_workbook

            wb = load_workbook(file_path, data_only=True, read_only=True)

            sheets_data = []
            total_rows = 0

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                rows = []
                for i, row in enumerate(sheet.iter_rows(values_only=True)):
                    if i >= self.MAX_TOTAL_ROWS:
                        break
                    rows.append(row)

                if not rows:
                    continue

                # 处理表头
                headers = [
                    str(h) if h else f"列{i + 1}" for i, h in enumerate(rows[0])
                ]
                data_rows = rows[1 : self.MAX_SAMPLE_ROWS + 1]
                row_count = len(rows) - 1
                total_rows += row_count

                # 推断数据类型
                data_types = self._infer_column_types(headers, data_rows)

                sheets_data.append(
                    {
                        "name": sheet_name,
                        "headers": headers,
                        "row_count": row_count,
                        "sample_rows": [
                            {
                                headers[i]: self._format_cell(cell)
                                for i, cell in enumerate(row)
                                if i < len(headers)
                            }
                            for row in data_rows
                        ],
                        "data_types": data_types,
                        "truncated": row_count > self.MAX_SAMPLE_ROWS,
                    }
                )

            wb.close()

            structured = {
                "type": "excel",
                "sheets": sheets_data,
                "total_rows": total_rows,
                "sheet_count": len(sheets_data),
            }

            text = self._to_text(structured)

            return ExtractedContent(
                raw_text=text,
                structured_data=structured,
                metadata={
                    "format": "excel",
                    "sheet_count": len(sheets_data),
                    "total_rows": total_rows,
                },
            )
        except ImportError:
            return ExtractedContent(
                raw_text="[需要安装 openpyxl 库来读取 Excel 文件]",
                structured_data=None,
                metadata={"error": "openpyxl not installed"},
            )
        except Exception as e:
            return ExtractedContent(
                raw_text=f"[Excel 读取失败: {e}]",
                structured_data=None,
                metadata={"error": str(e)},
            )

    def _format_cell(self, cell) -> str:
        """格式化单元格值"""
        if cell is None:
            return ""
        if isinstance(cell, (int, float)):
            return str(cell)
        return str(cell)

    def _infer_column_types(self, headers: list, rows: list) -> dict:
        """推断列数据类型"""
        types = {}
        for i, header in enumerate(headers):
            values = [row[i] for row in rows if i < len(row) and row[i] is not None]
            if not values:
                types[header] = "unknown"
            elif all(isinstance(v, (int, float)) for v in values):
                types[header] = "number"
            elif all(self._is_date(v) for v in values):
                types[header] = "date"
            else:
                types[header] = "string"
        return types

    def _is_date(self, value) -> bool:
        """检查是否为日期类型"""
        from datetime import datetime

        return isinstance(value, datetime)

    def _to_text(self, structured: dict) -> str:
        """将结构化数据转为可读文本"""
        lines = []
        for sheet in structured["sheets"]:
            lines.append(f"## Sheet: {sheet['name']} ({sheet['row_count']} 行)")
            lines.append("列: " + ", ".join(sheet["headers"]))
            lines.append("数据类型: " + json.dumps(sheet["data_types"], ensure_ascii=False))
            lines.append("")
            lines.append("示例数据:")
            for row in sheet["sample_rows"][:5]:
                lines.append("  " + json.dumps(row, ensure_ascii=False))
            if sheet["truncated"]:
                lines.append(f"  ... [共 {sheet['row_count']} 行] ...")
            lines.append("")
        return "\n".join(lines)


class PdfExtractor(BaseExtractor):
    """PDF 文件提取器"""

    MAX_PAGES = 20  # 最大页数

    def extract(self, file_path: Path) -> ExtractedContent:
        try:
            from pypdf import PdfReader

            reader = PdfReader(file_path)
            total_pages = len(reader.pages)

            sections = []
            text_parts = []

            for i, page in enumerate(reader.pages[: self.MAX_PAGES]):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    sections.append(
                        {"page": i + 1, "content": page_text[:2000]}  # 每页限制 2000 字符
                    )
                    text_parts.append(f"--- Page {i + 1} ---\n{page_text}")

            truncated = total_pages > self.MAX_PAGES

            structured = {
                "type": "document",
                "format": "pdf",
                "sections": sections,
                "total_pages": total_pages,
                "extracted_pages": min(total_pages, self.MAX_PAGES),
                "truncated": truncated,
            }

            text = "\n\n".join(text_parts)
            if truncated:
                text += f"\n\n... [共 {total_pages} 页，已提取前 {self.MAX_PAGES} 页] ..."

            return ExtractedContent(
                raw_text=text,
                structured_data=structured,
                metadata={
                    "format": "pdf",
                    "total_pages": total_pages,
                    "truncated": truncated,
                },
            )
        except ImportError:
            return ExtractedContent(
                raw_text="[需要安装 pypdf 库来读取 PDF 文件]",
                structured_data=None,
                metadata={"error": "pypdf not installed"},
            )
        except Exception as e:
            return ExtractedContent(
                raw_text=f"[PDF 读取失败: {e}]",
                structured_data=None,
                metadata={"error": str(e)},
            )


class DocxExtractor(BaseExtractor):
    """Word 文档提取器 (.docx)"""

    MAX_PARAGRAPHS = 500  # 最大段落数

    def extract(self, file_path: Path) -> ExtractedContent:
        try:
            from docx import Document

            doc = Document(str(file_path))

            paragraphs = []
            sections = []
            current_section = {"title": "正文", "content": [], "page": 1}

            for i, para in enumerate(doc.paragraphs[: self.MAX_PARAGRAPHS]):
                text = para.text.strip()
                if not text:
                    continue

                # 尝试识别标题
                if para.style and para.style.name and "Heading" in para.style.name:
                    # 保存当前 section
                    if current_section["content"]:
                        current_section["content"] = "\n".join(current_section["content"])
                        sections.append(current_section)
                    # 开始新 section
                    current_section = {"title": text, "content": [], "page": len(sections) + 1}
                else:
                    current_section["content"].append(text)

                paragraphs.append(text)

            # 保存最后一个 section
            if current_section["content"]:
                current_section["content"] = "\n".join(current_section["content"])
                sections.append(current_section)

            truncated = len(doc.paragraphs) > self.MAX_PARAGRAPHS

            structured = {
                "type": "document",
                "format": "docx",
                "sections": sections,
                "total_paragraphs": len(doc.paragraphs),
                "truncated": truncated,
            }

            text = "\n\n".join(paragraphs)
            if truncated:
                text += f"\n\n... [共 {len(doc.paragraphs)} 段，已提取前 {self.MAX_PARAGRAPHS} 段] ..."

            return ExtractedContent(
                raw_text=text,
                structured_data=structured,
                metadata={
                    "format": "docx",
                    "paragraph_count": len(doc.paragraphs),
                    "truncated": truncated,
                },
            )
        except ImportError:
            return ExtractedContent(
                raw_text="[需要安装 python-docx 库来读取 Word 文件]",
                structured_data=None,
                metadata={"error": "python-docx not installed"},
            )
        except Exception as e:
            return ExtractedContent(
                raw_text=f"[Word 文档读取失败: {e}]",
                structured_data=None,
                metadata={"error": str(e)},
            )


class PptxExtractor(BaseExtractor):
    """PowerPoint 文档提取器 (.pptx)"""

    MAX_SLIDES = 50  # 最大幻灯片数

    def extract(self, file_path: Path) -> ExtractedContent:
        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))

            slides_data = []
            text_parts = []

            slides_list = list(prs.slides)[:self.MAX_SLIDES]
            for i, slide in enumerate(slides_list):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = getattr(shape, "text", "")
                        if text and text.strip():
                            slide_text.append(text.strip())

                content = "\n".join(slide_text)
                slides_data.append({"slide": i + 1, "content": content[:1000]})  # 每页限制 1000 字符
                text_parts.append(f"--- Slide {i + 1} ---\n{content}")

            truncated = len(prs.slides) > self.MAX_SLIDES

            structured = {
                "type": "document",
                "format": "pptx",
                "sections": slides_data,
                "total_slides": len(prs.slides),
                "truncated": truncated,
            }

            text = "\n\n".join(text_parts)
            if truncated:
                text += f"\n\n... [共 {len(prs.slides)} 页，已提取前 {self.MAX_SLIDES} 页] ..."

            return ExtractedContent(
                raw_text=text,
                structured_data=structured,
                metadata={
                    "format": "pptx",
                    "slide_count": len(prs.slides),
                    "truncated": truncated,
                },
            )
        except ImportError:
            return ExtractedContent(
                raw_text="[需要安装 python-pptx 库来读取 PowerPoint 文件]",
                structured_data=None,
                metadata={"error": "python-pptx not installed"},
            )
        except Exception as e:
            return ExtractedContent(
                raw_text=f"[PowerPoint 读取失败: {e}]",
                structured_data=None,
                metadata={"error": str(e)},
            )


class FileExtractor:
    """统一文件提取器"""

    # 文件类型到提取器的映射
    _extractors: dict[str, BaseExtractor] = {
        # 文本文件
        ".txt": TextExtractor(),
        ".md": TextExtractor(),
        ".json": TextExtractor(),
        ".yaml": TextExtractor(),
        ".yml": TextExtractor(),
        ".xml": TextExtractor(),
        ".html": TextExtractor(),
        # CSV
        ".csv": CsvExtractor(),
        # Excel
        ".xlsx": ExcelExtractor(),
        ".xls": ExcelExtractor(),
        # PDF
        ".pdf": PdfExtractor(),
        # Word
        ".docx": DocxExtractor(),
        # PowerPoint
        ".pptx": PptxExtractor(),
        # 代码文件
        ".py": TextExtractor(),
        ".js": TextExtractor(),
        ".ts": TextExtractor(),
        ".tsx": TextExtractor(),
        ".jsx": TextExtractor(),
        ".java": TextExtractor(),
        ".cpp": TextExtractor(),
        ".c": TextExtractor(),
        ".h": TextExtractor(),
        ".go": TextExtractor(),
        ".rs": TextExtractor(),
        ".rb": TextExtractor(),
        ".php": TextExtractor(),
        ".sql": TextExtractor(),
        ".sh": TextExtractor(),
        ".bash": TextExtractor(),
        ".zsh": TextExtractor(),
        ".css": TextExtractor(),
        ".scss": TextExtractor(),
        ".less": TextExtractor(),
    }

    def extract(self, file_path: Path) -> ExtractedContent:
        """根据文件类型选择提取器并提取内容"""
        if not file_path.exists():
            return ExtractedContent(
                raw_text=f"[文件不存在: {file_path}]",
                structured_data=None,
                metadata={"error": "file not found"},
            )

        suffix = file_path.suffix.lower()
        extractor = self._extractors.get(suffix, TextExtractor())
        return extractor.extract(file_path)

    def get_supported_extensions(self) -> list[str]:
        """获取支持的文件扩展名列表"""
        return list(self._extractors.keys())

    def is_supported(self, file_path: Path) -> bool:
        """检查文件类型是否支持"""
        return file_path.suffix.lower() in self._extractors
