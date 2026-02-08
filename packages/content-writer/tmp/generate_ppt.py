#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# 设置颜色方案
PRIMARY_COLOR = RGBColor(59, 130, 246)  # 蓝色
SECONDARY_COLOR = RGBColor(51, 51, 51)  # 深灰
ACCENT_COLOR = RGBColor(16, 185, 129)  # 绿色

# 幻灯片1：封面
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
left = Inches(1)
top = Inches(2.5)
width = Inches(8)
height = Inches(2)

# 标题
title_box = slide1.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_frame.text = "UI设计技术栈调研"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(54)
title_para.font.bold = True
title_para.font.color.rgb = PRIMARY_COLOR
title_para.alignment = PP_ALIGN.CENTER

# 副标题
subtitle_box = slide1.shapes.add_textbox(left, Inches(4.2), width, Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "2025年核心趋势与技术选型指南"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(24)
subtitle_para.font.color.rgb = SECONDARY_COLOR
subtitle_para.alignment = PP_ALIGN.CENTER

# 幻灯片2：要点1 - AI驱动设计
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame2 = title_box2.text_frame
title_frame2.text = "核心要点 1：AI驱动设计成为主流"
title_para2 = title_frame2.paragraphs[0]
title_para2.font.size = Pt(36)
title_para2.font.bold = True
title_para2.font.color.rgb = PRIMARY_COLOR

# 内容框
content_left = Inches(1)
content_top = Inches(1.8)
content_width = Inches(8)
content_height = Inches(4.5)

content_box2 = slide2.shapes.add_textbox(content_left, content_top, content_width, content_height)
tf2 = content_box2.text_frame
tf2.word_wrap = True

# 要点1
p1 = tf2.paragraphs[0]
p1.text = "🚀 v0.dev + shadcn/ui 组合"
p1.font.size = Pt(24)
p1.font.bold = True
p1.font.color.rgb = SECONDARY_COLOR
p1.space_after = Pt(10)

p1_sub = tf2.add_paragraph()
p1_sub.text = "自然语言生成生产级React代码，与Next.js深度集成"
p1_sub.font.size = Pt(18)
p1_sub.font.color.rgb = RGBColor(100, 100, 100)
p1_sub.level = 1
p1_sub.space_after = Pt(20)

# 要点2
p2 = tf2.add_paragraph()
p2.text = "🎨 设计工具智能化"
p2.font.size = Pt(24)
p2.font.bold = True
p2.font.color.rgb = SECONDARY_COLOR
p2.space_after = Pt(10)

p2_sub = tf2.add_paragraph()
p2_sub.text = "Figma Make、Galileo AI等工具打破设计开发边界"
p2_sub.font.size = Pt(18)
p2_sub.font.color.rgb = RGBColor(100, 100, 100)
p2_sub.level = 1
p2_sub.space_after = Pt(20)

# 要点3
p3 = tf2.add_paragraph()
p3.text = "⚡ 效率提升显著"
p3.font.size = Pt(24)
p3.font.bold = True
p3.font.color.rgb = SECONDARY_COLOR
p3.space_after = Pt(10)

p3_sub = tf2.add_paragraph()
p3_sub.text = "AI工具可将设计开发效率提升10倍以上"
p3_sub.font.size = Pt(18)
p3_sub.font.color.rgb = RGBColor(100, 100, 100)
p3_sub.level = 1

# 幻灯片3：要点2 - shadcn/ui
slide3 = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title_box3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame3 = title_box3.text_frame
title_frame3.text = "核心要点 2：shadcn/ui引领组件库新范式"
title_para3 = title_frame3.paragraphs[0]
title_para3.font.size = Pt(36)
title_para3.font.bold = True
title_para3.font.color.rgb = PRIMARY_COLOR

# 内容
content_box3 = slide3.shapes.add_textbox(content_left, content_top, content_width, content_height)
tf3 = content_box3.text_frame
tf3.word_wrap = True

# 要点1
p1 = tf3.paragraphs[0]
p1.text = "⭐ 85k+ GitHub Stars"
p1.font.size = Pt(24)
p1.font.bold = True
p1.font.color.rgb = SECONDARY_COLOR
p1.space_after = Pt(10)

p1_sub = tf3.add_paragraph()
p1_sub.text = "不是传统组件库，而是「组件分发系统」"
p1_sub.font.size = Pt(18)
p1_sub.font.color.rgb = RGBColor(100, 100, 100)
p1_sub.level = 1
p1_sub.space_after = Pt(20)

# 要点2
p2 = tf3.add_paragraph()
p2.text = "🔓 完全开放可控"
p2.font.size = Pt(24)
p2.font.bold = True
p2.font.color.rgb = SECONDARY_COLOR
p2.space_after = Pt(10)

p2_sub = tf3.add_paragraph()
p2_sub.text = "代码完全开放，无黑盒，专为AI工具优化"
p2_sub.font.size = Pt(18)
p2_sub.font.color.rgb = RGBColor(100, 100, 100)
p2_sub.level = 1
p2_sub.space_after = Pt(20)

# 要点3
p3 = tf3.add_paragraph()
p3.text = "🏢 顶级公司采用"
p3.font.size = Pt(24)
p3.font.bold = True
p3.font.color.rgb = SECONDARY_COLOR
p3.space_after = Pt(10)

p3_sub = tf3.add_paragraph()
p3_sub.text = "OpenAI、Adobe、Sonos等公司生产环境使用"
p3_sub.font.size = Pt(18)
p3_sub.font.color.rgb = RGBColor(100, 100, 100)
p3_sub.level = 1

# 幻灯片4：要点3 - 推荐技术栈
slide4 = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title_box4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame4 = title_box4.text_frame
title_frame4.text = "核心要点 3：推荐技术栈组合"
title_para4 = title_frame4.paragraphs[0]
title_para4.font.size = Pt(36)
title_para4.font.bold = True
title_para4.font.color.rgb = PRIMARY_COLOR

# 内容
content_box4 = slide4.shapes.add_textbox(content_left, content_top, content_width, content_height)
tf4 = content_box4.text_frame
tf4.word_wrap = True

# 方案A
p1 = tf4.paragraphs[0]
p1.text = "🚀 方案A：快速开发（推荐）"
p1.font.size = Pt(24)
p1.font.bold = True
p1.font.color.rgb = ACCENT_COLOR
p1.space_after = Pt(10)

p1_sub = tf4.add_paragraph()
p1_sub.text = "Figma + shadcn/ui + Tailwind CSS + v0.dev"
p1_sub.font.size = Pt(18)
p1_sub.font.color.rgb = RGBColor(100, 100, 100)
p1_sub.level = 1
p1_sub.space_after = Pt(20)

# 方案B
p2 = tf4.add_paragraph()
p2.text = "🏢 方案B：企业级应用"
p2.font.size = Pt(24)
p2.font.bold = True
p2.font.color.rgb = SECONDARY_COLOR
p2.space_after = Pt(10)

p2_sub = tf4.add_paragraph()
p2_sub.text = "Figma + Material UI + Emotion"
p2_sub.font.size = Pt(18)
p2_sub.font.color.rgb = RGBColor(100, 100, 100)
p2_sub.level = 1
p2_sub.space_after = Pt(20)

# 方案C
p3 = tf4.add_paragraph()
p3.text = "🎨 方案C：高度定制"
p3.font.size = Pt(24)
p3.font.bold = True
p3.font.color.rgb = SECONDARY_COLOR
p3.space_after = Pt(10)

p3_sub = tf4.add_paragraph()
p3_sub.text = "Figma + Radix UI + Tailwind CSS（完全掌控样式）"
p3_sub.font.size = Pt(18)
p3_sub.font.color.rgb = RGBColor(100, 100, 100)
p3_sub.level = 1

# 幻灯片5：总结
slide5 = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title_box5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame5 = title_box5.text_frame
title_frame5.text = "总结：2025年UI设计三大趋势"
title_para5 = title_frame5.paragraphs[0]
title_para5.font.size = Pt(36)
title_para5.font.bold = True
title_para5.font.color.rgb = PRIMARY_COLOR

# 内容
content_box5 = slide5.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4))
tf5 = content_box5.text_frame
tf5.word_wrap = True

# 趋势1
p1 = tf5.paragraphs[0]
p1.text = "1️⃣  AI 驱动"
p1.font.size = Pt(28)
p1.font.bold = True
p1.font.color.rgb = PRIMARY_COLOR
p1.space_after = Pt(15)

# 趋势2
p2 = tf5.add_paragraph()
p2.text = "2️⃣  开放代码"
p2.font.size = Pt(28)
p2.font.bold = True
p2.font.color.rgb = PRIMARY_COLOR
p2.space_after = Pt(15)

# 趋势3
p3 = tf5.add_paragraph()
p3.text = "3️⃣  设计开发一体化"
p3.font.size = Pt(28)
p3.font.bold = True
p3.font.color.rgb = PRIMARY_COLOR
p3.space_after = Pt(30)

# 结束语
p4 = tf5.add_paragraph()
p4.text = "优先考虑 shadcn/ui + v0.dev 组合"
p4.font.size = Pt(22)
p4.font.color.rgb = ACCENT_COLOR
p4.alignment = PP_ALIGN.CENTER

# 保存文件
output_path = "/tmp/UI设计技术栈调研-核心要点.pptx"
prs.save(output_path)
print(f"✅ PPT已生成：{output_path}")
print(f"📊 共 {len(prs.slides)} 张幻灯片")
